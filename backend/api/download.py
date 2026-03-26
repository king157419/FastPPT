import os
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse, JSONResponse
import base64

router = APIRouter()

OUTPUT_DIR = "outputs"


@router.get("/download/{filename}")
def download_file(filename: str):
    path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(path):
        raise HTTPException(404, "文件不存在")
    return FileResponse(path, filename=filename)


@router.get("/preview/{filename}")
def preview_pptx(filename: str):
    """返回 PPT 各页缩略图（base64 PNG 列表）"""
    path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(path):
        raise HTTPException(404, "文件不存在")

    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        import io

        # 加载中文字体，fallback 顺序
        def _load_font(size: int):
            for path in [
                "C:/Windows/Fonts/msyh.ttc",
                "C:/Windows/Fonts/simhei.ttf",
                "C:/Windows/Fonts/simsun.ttc",
            ]:
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    continue
            return ImageFont.load_default()

        font_title = _load_font(30)
        font_body = _load_font(19)
        font_page = _load_font(15)

        prs = Presentation(path)
        slides_data = []

        for i, slide in enumerate(prs.slides):
            # 创建缩略图画布 (960x540)
            img = Image.new("RGB", (960, 540), color=(26, 26, 46))
            draw = ImageDraw.Draw(img)

            # 提取幻灯片中的文字
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            texts.append(t)

            # 顶部金色装饰条 + 页码
            draw.rectangle([0, 0, 960, 6], fill=(255, 215, 0))
            draw.text((20, 16), f"第 {i+1} 页", fill=(180, 196, 222), font=font_page)

            # 绘制文字内容
            y = 56
            for j, text in enumerate(texts[:8]):
                if y > 490:
                    break
                snippet = text[:50] + "..." if len(text) > 50 else text
                if j == 0:
                    draw.text((40, y), snippet, fill=(255, 215, 0), font=font_title)
                    y += 42
                else:
                    draw.text((48, y), f"• {snippet}", fill=(210, 220, 240), font=font_body)
                    y += 30

            # 转 base64
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode()
            slides_data.append({"page": i + 1, "image": f"data:image/png;base64,{b64}"})

        return JSONResponse({"slides": slides_data, "total": len(slides_data)})

    except Exception as e:
        raise HTTPException(500, f"预览生成失败：{e}")
