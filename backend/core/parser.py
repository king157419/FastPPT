"""Document parsing utilities with structured chunk metadata."""
from __future__ import annotations

from pathlib import Path


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        rows = _parse_pdf_rows(file_path)
        return "\n".join(item["text"] for item in rows)
    if ext in (".docx", ".doc"):
        rows = _parse_docx_rows(file_path)
        return "\n".join(item["text"] for item in rows)
    if ext in (".pptx", ".ppt"):
        rows = _parse_pptx_rows(file_path)
        return "\n".join(item["text"] for item in rows)
    if ext in (".txt", ".md"):
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return ""


def extract_ppt_outline(file_path: str, max_items: int = 20) -> list[str]:
    rows = _parse_pptx_rows(file_path)
    outline: list[str] = []
    for row in rows:
        title = str(row.get("title", "")).strip() or str(row.get("text", "")).split("\n")[0].strip()
        if not title:
            continue
        outline.append(title)
        if len(outline) >= max_items:
            break
    return outline


def extract_chunks_with_metadata(file_path: str) -> list[dict]:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        rows = _parse_pdf_rows(file_path)
        chunks: list[dict] = []
        for idx, row in enumerate(rows):
            text = row["text"].strip()
            if not text:
                continue
            chunks.append(
                {
                    "text": text,
                    "metadata": {
                        "page_or_slide": f"page-{row['page']}",
                        "chunk_index": idx,
                    },
                }
            )
        return chunks

    if ext in (".docx", ".doc"):
        rows = _parse_docx_rows(file_path)
        chunks: list[dict] = []
        for idx, row in enumerate(rows):
            text = row["text"].strip()
            if not text:
                continue
            section = row.get("section") or f"section-{idx + 1}"
            chunks.append(
                {
                    "text": text,
                    "metadata": {
                        "page_or_slide": section,
                        "chunk_index": idx,
                    },
                }
            )
        return chunks

    if ext in (".pptx", ".ppt"):
        rows = _parse_pptx_rows(file_path)
        chunks: list[dict] = []
        for idx, row in enumerate(rows):
            text = row["text"].strip()
            if not text:
                continue
            chunks.append(
                {
                    "text": text,
                    "metadata": {
                        "page_or_slide": f"slide-{row['slide']}",
                        "chunk_index": idx,
                    },
                }
            )
        return chunks

    text = extract_text(file_path)
    plain_chunks = chunk_text(text)
    return [
        {
            "text": item,
            "metadata": {"page_or_slide": "text", "chunk_index": idx},
        }
        for idx, item in enumerate(plain_chunks)
        if item.strip()
    ]


def _parse_pdf_rows(path: str) -> list[dict]:
    import pdfplumber

    rows: list[dict] = []
    with pdfplumber.open(path) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                rows.append({"page": page_idx, "text": text})
    return rows


def _parse_docx_rows(path: str) -> list[dict]:
    from docx import Document

    doc = Document(path)
    rows: list[dict] = []
    current_heading = "section-1"
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue

        style_name = (para.style.name or "") if para.style else ""
        if style_name.lower().startswith("heading"):
            current_heading = text
            continue

        rows.append({"section": current_heading, "text": text})
    return rows


def _parse_pptx_rows(path: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(path)
    rows: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)

        if not lines:
            continue

        rows.append(
            {
                "slide": slide_idx,
                "title": lines[0],
                "text": "\n".join(lines),
            }
        )
    return rows


def chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> list[str]:
    """Split text by character length with overlap."""
    chunks = []
    cursor = 0
    body = text or ""
    step = max(1, chunk_size - overlap)
    while cursor < len(body):
        chunks.append(body[cursor : cursor + chunk_size])
        cursor += step
    return chunks
