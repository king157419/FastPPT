"""使用 python-pptx 生成 .pptx 课件"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import time
import logging
import requests
from typing import Optional

logger = logging.getLogger(__name__)


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

COLOR_BG = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝背景
COLOR_ACCENT = RGBColor(0x16, 0x21, 0x3E)   # 次背景
COLOR_HL = RGBColor(0x0F, 0x3A, 0x7A)       # 蓝色高亮
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
COLOR_LIGHT = RGBColor(0xB0, 0xC4, 0xDE)


def _add_bg(slide, prs):
    """给幻灯片填充深色背景"""
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def _add_textbox(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or COLOR_WHITE
    return txBox


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _cover_slide(prs, topic, audience, duration):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _add_bg(slide, prs)
    # 顶部装饰条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), COLOR_YELLOW)
    # 主标题
    _add_textbox(slide, topic,
                 Inches(1), Inches(2), Inches(11), Inches(1.5),
                 font_size=44, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # 副标题
    sub = f"适用对象：{audience}　|　课时：{duration}"
    _add_textbox(slide, sub,
                 Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                 font_size=20, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    # 底部
    _add_textbox(slide, "TeachMind · AI 智能备课系统",
                 Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 font_size=14, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    return slide


def _toc_slide(prs, key_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), COLOR_YELLOW)
    _add_textbox(slide, "目  录",
                 Inches(0.8), Inches(0.3), Inches(6), Inches(0.8),
                 font_size=32, bold=True, color=COLOR_YELLOW)
    for i, kp in enumerate(key_points):
        y = Inches(1.4 + i * 0.85)
        _add_rect(slide, Inches(0.8), y + Inches(0.1), Inches(0.35), Inches(0.45), COLOR_HL)
        _add_textbox(slide, str(i + 1),
                     Inches(0.85), y, Inches(0.3), Inches(0.65),
                     font_size=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(slide, kp,
                     Inches(1.4), y, Inches(10), Inches(0.65),
                     font_size=22, color=COLOR_WHITE)
    return slide


def _content_slide(prs, title, bullets, rag_context="", tip=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, COLOR_HL)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), COLOR_YELLOW)
    _add_textbox(slide, title,
                 Inches(0.8), Inches(0.2), Inches(11), Inches(0.9),
                 font_size=30, bold=True, color=COLOR_YELLOW)
    # 分隔线效果
    _add_rect(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.04), COLOR_HL)
    # 内容要点
    for i, b in enumerate(bullets):
        y = Inches(1.35 + i * 0.85)
        if y + Inches(0.85) > SLIDE_H - Inches(0.3):
            break
        _add_rect(slide, Inches(0.8), y + Inches(0.18), Inches(0.12), Inches(0.35), COLOR_YELLOW)
        _add_textbox(slide, b,
                     Inches(1.1), y, Inches(11.2), Inches(0.75),
                     font_size=20, color=COLOR_WHITE)
    # 教师提示（右下角黄色小字）
    if tip:
        _add_textbox(slide, f"💡 {tip}",
                     Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.45),
                     font_size=12, color=COLOR_YELLOW)
    elif rag_context:
        snippet = rag_context[:80] + "..." if len(rag_context) > 80 else rag_context
        _add_textbox(slide, f"📚 {snippet}",
                     Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.45),
                     font_size=11, color=COLOR_LIGHT)
    return slide


def _summary_slide(prs, topic, key_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), COLOR_YELLOW)
    _add_textbox(slide, "课程小结",
                 Inches(0.8), Inches(0.3), Inches(6), Inches(0.8),
                 font_size=32, bold=True, color=COLOR_YELLOW)
    _add_textbox(slide, f"本节课围绕「{topic}」，共覆盖以下核心知识点：",
                 Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
                 font_size=20, color=COLOR_LIGHT)
    for i, kp in enumerate(key_points):
        y = Inches(2.1 + i * 0.75)
        _add_textbox(slide, f"✓  {kp}",
                     Inches(1.2), y, Inches(10), Inches(0.65),
                     font_size=20, color=COLOR_WHITE)
    _add_textbox(slide, "感谢聆听，欢迎提问！",
                 Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                 font_size=22, bold=True, color=COLOR_YELLOW, align=PP_ALIGN.CENTER)
    return slide


def _convert_to_pptxgenjs_format(intent: dict, slide_contents: list[dict]) -> dict:
    """
    将现有的 intent 和 slide_contents 转换为 PptxGenJS 服务所需的格式。
    """
    topic = intent.get("topic", "未知主题")
    audience = intent.get("audience", "学生")
    key_points = intent.get("key_points", ["核心概念", "基本原理", "实际应用"])
    duration = intent.get("duration", "45分钟")

    pages = []

    # Cover slide
    pages.append({
        "type": "cover",
        "title": topic,
        "subtitle": f"适用对象：{audience}　|　课时：{duration}"
    })

    # Table of contents
    pages.append({
        "type": "agenda",
        "title": "目  录",
        "items": key_points
    })

    # Content slides
    for i, sc in enumerate(slide_contents):
        kp = sc.get("key_point", key_points[i] if i < len(key_points) else f"知识点{i+1}")
        bullets = sc.get("bullets", [f"{kp}的核心内容"])
        tip = sc.get("tip", "")

        pages.append({
            "type": "content",
            "title": f"{i+1}. {kp}",
            "bullets": bullets,
            "tip": tip
        })

    # Summary slide
    pages.append({
        "type": "summary",
        "title": "课程小结",
        "takeaways": key_points
    })

    return {
        "teaching_spec": {
            "title": topic,
            "subject": intent.get("subject", ""),
            "teacher": intent.get("teacher", "TeachMind AI")
        },
        "pages": pages,
        "theme": "education",
        "metadata": {
            "author": intent.get("teacher", "TeachMind AI"),
            "title": topic,
            "filename": os.path.basename(output_path) if 'output_path' in locals() else f"{topic}.pptx"
        }
    }


def call_pptxgenjs_service(
    intent: dict,
    slide_contents: list[dict],
    output_path: str,
    service_url: Optional[str] = None,
    timeout: int = 30,
    max_retries: int = 2
) -> Optional[str]:
    """
    调用 PptxGenJS 服务生成 PPT。

    Args:
        intent: 课程意图字典
        slide_contents: 幻灯片内容列表
        output_path: 输出文件路径
        service_url: PptxGenJS 服务地址，默认从环境变量读取
        timeout: 请求超时时间（秒）
        max_retries: 最大重试次数

    Returns:
        成功返回输出文件路径，失败返回 None
    """
    if service_url is None:
        service_url = os.environ.get("PPTXGENJS_SERVICE_URL", "http://localhost:3000")

    endpoint = f"{service_url}/generate"

    # 转换数据格式
    payload = _convert_to_pptxgenjs_format(intent, slide_contents)
    payload["metadata"]["filename"] = os.path.basename(output_path)

    start_time = time.time()

    for attempt in range(max_retries + 1):
        try:
            logger.info(f"[PptxGenJS] 调用服务生成 PPT (尝试 {attempt + 1}/{max_retries + 1})")

            response = requests.post(
                endpoint,
                json=payload,
                timeout=timeout,
                headers={"Content-Type": "application/json"}
            )

            if response.status_code == 200:
                # 保存返回的 PPTX 文件
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                with open(output_path, "wb") as f:
                    f.write(response.content)

                duration = time.time() - start_time
                slide_count = len(payload["pages"])
                avg_time = duration / slide_count if slide_count > 0 else 0

                logger.info(
                    f"[PptxGenJS] 生成成功: {slide_count} 页, "
                    f"耗时 {duration:.2f}s ({avg_time:.2f}s/页)"
                )

                return output_path
            else:
                logger.warning(
                    f"[PptxGenJS] 服务返回错误: {response.status_code} - {response.text[:200]}"
                )

        except requests.exceptions.Timeout:
            logger.warning(f"[PptxGenJS] 请求超时 (尝试 {attempt + 1}/{max_retries + 1})")
        except requests.exceptions.ConnectionError:
            logger.warning(f"[PptxGenJS] 连接失败 (尝试 {attempt + 1}/{max_retries + 1})")
        except Exception as e:
            logger.error(f"[PptxGenJS] 调用失败: {e}")
            break

        # 重试前等待
        if attempt < max_retries:
            time.sleep(1)

    duration = time.time() - start_time
    logger.warning(f"[PptxGenJS] 所有尝试失败，耗时 {duration:.2f}s")
    return None


def _generate_pptx_pythonpptx(intent: dict, slide_contents: list[dict], output_path: str) -> str:
    """
    使用 python-pptx 生成 PPT（fallback 方法）。
    """
    topic = intent.get("topic", "未知主题")
    audience = intent.get("audience", "学生")
    key_points = intent.get("key_points", ["核心概念", "基本原理", "实际应用"])
    duration = intent.get("duration", "45分钟")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _cover_slide(prs, topic, audience, duration)
    _toc_slide(prs, key_points)

    for i, sc in enumerate(slide_contents):
        kp = sc.get("key_point", key_points[i] if i < len(key_points) else f"知识点{i+1}")
        bullets = sc.get("bullets", [f"{kp}的核心内容"])
        tip = sc.get("tip", "")
        rag_ctx = sc.get("rag_ctx", "")
        _content_slide(prs, f"{i+1}. {kp}", bullets, rag_ctx, tip)

    _summary_slide(prs, topic, key_points)

    prs.save(output_path)
    return output_path


def generate_pptx(intent: dict, slide_contents: list[dict], output_path: str) -> str:
    """
    根据 intent JSON 和 slide_contents 生成 .pptx 文件。
    优先使用 PptxGenJS 服务，失败时 fallback 到 python-pptx。

    Args:
        intent: 课程意图字典
        slide_contents: [{"key_point": str, "bullets": [str], "tip": str, "rag_ctx": str}]
        output_path: 输出文件路径

    Returns:
        输出文件路径
    """
    start_time = time.time()

    # 尝试使用 PptxGenJS 服务
    result = call_pptxgenjs_service(intent, slide_contents, output_path)

    if result:
        duration = time.time() - start_time
        logger.info(f"[PPT生成] 使用 PptxGenJS 服务成功，总耗时 {duration:.2f}s")
        return result

    # Fallback 到 python-pptx
    logger.info("[PPT生成] PptxGenJS 服务不可用，使用 python-pptx fallback")
    fallback_start = time.time()

    result = _generate_pptx_pythonpptx(intent, slide_contents, output_path)

    fallback_duration = time.time() - fallback_start
    total_duration = time.time() - start_time
    logger.info(
        f"[PPT生成] python-pptx fallback 完成，"
        f"生成耗时 {fallback_duration:.2f}s，总耗时 {total_duration:.2f}s"
    )

    return result
