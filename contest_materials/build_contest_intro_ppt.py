from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


W = Inches(13.333)
H = Inches(7.5)


def add_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, text, x=0.8, y=0.45, size=42, color=RGBColor(255, 255, 255)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(11.7), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = color


def add_subtitle(slide, text, x=0.82, y=1.15, size=18, color=RGBColor(203, 213, 225)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(11.6), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = color


def add_footer(slide, text, color=RGBColor(148, 163, 184)):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(5.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = color


def add_card(slide, x, y, w, h, title, lines, fill, title_color=RGBColor(255, 255, 255), body_color=RGBColor(241, 245, 249)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.16), Inches(w - 0.44), Inches(0.45))
    p = tbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = title_color

    bbox = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.72), Inches(w - 0.56), Inches(h - 0.86))
    tf = bbox.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = f"• {line}"
        para.font.size = Pt(14)
        para.font.name = "Microsoft YaHei"
        para.font.color.rgb = body_color


def add_banner(slide, x, y, w, h, text, fill, text_color=RGBColor(15, 23, 42), size=18):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(h - 0.2))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.font.color.rgb = text_color


def add_stat_box(slide, x, y, w, h, top, bottom, fill):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()

    b1 = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.18), Inches(w - 0.36), Inches(0.45))
    p1 = b1.text_frame.paragraphs[0]
    p1.text = top
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.name = "Microsoft YaHei"
    p1.font.color.rgb = RGBColor(255, 255, 255)

    b2 = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.8), Inches(w - 0.36), Inches(0.75))
    p2 = b2.text_frame.paragraphs[0]
    p2.text = bottom
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.name = "Consolas"
    p2.font.color.rgb = RGBColor(255, 255, 255)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(7, 22, 51))
    tri = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE, Inches(10.9), Inches(0.0), Inches(2.4), Inches(2.4))
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(251, 146, 60)
    tri.line.fill.background()
    add_title(s, "FastPPT 项目简介", y=1.8, size=54)
    add_subtitle(s, "A04 多模态 AI 互动式教学智能体", y=3.0, size=24, color=RGBColor(191, 219, 254))
    add_subtitle(s, "不是普通 AI PPT 生成器，而是高校教师的教学设计执行系统", y=3.7, size=21, color=RGBColor(147, 197, 253))
    add_footer(s, "Contest Submission Deck")

    # 2 Understanding
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(250, 245, 238))
    add_title(s, "01 我们对赛题的理解", color=RGBColor(17, 24, 39))
    add_subtitle(s, "真正难的不是“生成”，而是“理解、利用资料、支持修改”", color=RGBColor(87, 83, 78))
    add_card(
        s,
        0.8,
        2.0,
        3.8,
        3.9,
        "不是",
        ["给一个主题就自动出整套 PPT", "只追求速度和视觉炫技", "生成后无法继续稳定修改"],
        RGBColor(127, 29, 29),
    )
    add_card(
        s,
        4.8,
        2.0,
        3.8,
        3.9,
        "而是",
        ["先理解教师真实教学意图", "再把旧 PPT 和资料真正用上", "最后输出可改、可导出的课件"],
        RGBColor(30, 64, 175),
    )
    add_card(
        s,
        8.8,
        2.0,
        3.7,
        3.9,
        "评审价值",
        ["多模态输入", "RAG 与证据绑定", "页级 revise 闭环", "教学实用性优先"],
        RGBColor(5, 150, 105),
    )
    add_footer(s, "Page 2", color=RGBColor(120, 113, 108))

    # 3 user and scenario
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(17, 24, 39))
    add_title(s, "02 目标用户与主场景")
    add_subtitle(s, "主战场不是从零生成，而是旧课增强与资料重构")
    add_card(
        s,
        0.9,
        1.9,
        3.6,
        4.5,
        "目标用户",
        ["有真实授课任务的高校教师", "已经具备教学判断和课程思路", "手里已有教材、旧课件、案例和素材"],
        RGBColor(37, 99, 235),
    )
    add_card(
        s,
        4.85,
        1.9,
        3.6,
        4.5,
        "主场景",
        ["基于旧 PPT 或同类课件", "补最新案例、数据、论文、政策", "增加更相关的示意图和图表"],
        RGBColor(217, 119, 6),
    )
    add_card(
        s,
        8.8,
        1.9,
        3.7,
        4.5,
        "老师最在意",
        ["别乱猜", "别太花", "图要相关", "修改时别带乱其他页"],
        RGBColor(8, 145, 178),
    )
    add_footer(s, "Page 3")

    # 4 solution chain
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(255, 255, 255))
    add_title(s, "03 方案主链路", color=RGBColor(17, 24, 39))
    add_subtitle(s, "把教学意图、资料、生成和修改串成连续闭环", color=RGBColor(71, 85, 105))
    steps = [
        ("Teacher Input", RGBColor(37, 99, 235)),
        ("TeachingSpec", RGBColor(22, 163, 74)),
        ("SourceAnchor", RGBColor(147, 51, 234)),
        ("Generate", RGBColor(234, 88, 12)),
        ("Revise / Export", RGBColor(15, 23, 42)),
    ]
    x = 0.75
    for idx, (txt, color) in enumerate(steps):
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.0), Inches(2.2), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        box = s.shapes.add_textbox(Inches(x), Inches(3.33), Inches(2.2), Inches(0.45))
        p = box.text_frame.paragraphs[0]
        p.text = txt
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Consolas" if idx == 0 else "Microsoft YaHei"
        p.font.color.rgb = RGBColor(255, 255, 255)
        if idx < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 2.28), Inches(3.2), Inches(0.55), Inches(0.68))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(203, 213, 225)
            arr.line.fill.background()
        x += 2.5
    add_banner(s, 1.15, 5.0, 11.0, 0.85, "生成前强确认：教学目标 / 面向学生类型 / 重点难点", RGBColor(254, 240, 138))
    add_footer(s, "Page 4", color=RGBColor(120, 130, 150))

    # 5 current abilities
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(9, 30, 66))
    add_title(s, "04 当前已具备的可演示能力")
    add_subtitle(s, "这部分强调“已跑通”，不把目标态写成已完成")
    add_card(
        s,
        0.85,
        1.95,
        3.7,
        4.2,
        "输入与理解",
        ["多文件上传", "对话采集教学意图", "TeachingSpec 编译与预检"],
        RGBColor(30, 64, 175),
    )
    add_card(
        s,
        4.8,
        1.95,
        3.7,
        4.2,
        "生成与预览",
        ["输出 slides_json", "前端预览课件", "支持真实 PPTX / DOCX 导出"],
        RGBColor(5, 150, 105),
    )
    add_card(
        s,
        8.75,
        1.95,
        3.7,
        4.2,
        "修改与演示",
        ["支持页级 revise", "可重新导出结果", "系统已运行在稳定 demo 模式"],
        RGBColor(217, 119, 6),
    )
    add_footer(s, "Page 5")

    # 6 evidence/status
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(245, 247, 250))
    add_title(s, "05 当前系统事实", color=RGBColor(17, 24, 39))
    add_subtitle(s, "代码与运行状态共同证明：项目已具备端到端演示基础", color=RGBColor(71, 85, 105))
    add_stat_box(s, 0.9, 2.0, 2.35, 1.8, "Health", "healthy", RGBColor(22, 163, 74))
    add_stat_box(s, 3.45, 2.0, 2.35, 1.8, "Mode", "demo_mode=true", RGBColor(37, 99, 235))
    add_stat_box(s, 6.0, 2.0, 2.35, 1.8, "Chat", "plain", RGBColor(124, 58, 237))
    add_stat_box(s, 8.55, 2.0, 2.35, 1.8, "RAG", "tfidf", RGBColor(8, 145, 178))
    add_stat_box(s, 11.1, 2.0, 1.3, 1.8, "Export", "PPTX", RGBColor(234, 88, 12))
    add_card(
        s,
        1.05,
        4.25,
        5.4,
        1.65,
        "代码骨架已落地",
        ["已有 TeachingSpec 对象", "已有 slide blocks 协议"],
        RGBColor(30, 41, 59),
    )
    add_card(
        s,
        6.75,
        4.25,
        5.5,
        1.65,
        "仍在推进",
        ["两阶段生成、blocks 全量渲染、课程记忆与 verify/repair 闭环"],
        RGBColor(100, 116, 139),
    )
    add_footer(s, "Page 6", color=RGBColor(120, 130, 150))

    # 7 innovation
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(24, 24, 27))
    add_title(s, "06 差异化创新点")
    add_subtitle(s, "关键不在模型更大，而在工程主链更适合教师工作流")
    add_card(
        s,
        0.85,
        1.95,
        2.8,
        4.35,
        "Spec First",
        ["先把需求收束成教学规格", "避免 prompt 驱动的随机漂移"],
        RGBColor(37, 99, 235),
    )
    add_card(
        s,
        3.85,
        1.95,
        2.8,
        4.35,
        "SourceAnchor",
        ["不是普通检索", "而是可绑定、可追溯的教学证据"],
        RGBColor(22, 163, 74),
    )
    add_card(
        s,
        6.85,
        1.95,
        2.8,
        4.35,
        "Page-level Revise",
        ["支持只改一页或补一个案例", "更贴近老师真实修改习惯"],
        RGBColor(217, 119, 6),
    )
    add_card(
        s,
        9.85,
        1.95,
        2.65,
        4.35,
        "CourseMemory",
        ["课程资产长期沉淀", "形成复利型工作区"],
        RGBColor(168, 85, 247),
    )
    add_footer(s, "Page 7")

    # 8 demo flow
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(255, 255, 255))
    add_title(s, "07 答辩演示路径", color=RGBColor(17, 24, 39))
    add_subtitle(s, "用最强主场景讲清楚价值：旧 PPT 增强，而不是一键出稿", color=RGBColor(71, 85, 105))
    add_card(
        s,
        0.9,
        2.0,
        11.6,
        3.7,
        "Demo Script",
        [
            "上传旧 PPT + 教材片段",
            "确认教学目标、学生类型、重点难点",
            "触发生成并查看 slides 预览",
            "对指定页提出修改要求",
            "导出 PPTX 与 DOCX 作为交付物",
        ],
        RGBColor(30, 41, 59),
    )
    add_banner(s, 1.15, 6.0, 5.3, 0.65, "Frontend: http://localhost:5173", RGBColor(219, 234, 254), RGBColor(30, 41, 59), 15)
    add_banner(s, 6.9, 6.0, 5.3, 0.65, "Backend: http://localhost:8000", RGBColor(220, 252, 231), RGBColor(22, 101, 52), 15)
    add_footer(s, "Page 8", color=RGBColor(120, 130, 150))

    # 9 roadmap
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(250, 245, 238))
    add_title(s, "08 路线图与当前边界", color=RGBColor(17, 24, 39))
    add_subtitle(s, "先把“可控首稿”做稳，再把“可用性”做强", color=RGBColor(120, 113, 108))
    add_card(
        s,
        0.85,
        2.0,
        3.8,
        4.1,
        "P0",
        ["TeachingSpec 继续强化", "SourceAnchor 最小闭环", "两阶段生成打通", "前端支持 blocks 渲染"],
        RGBColor(37, 99, 235),
    )
    add_card(
        s,
        4.8,
        2.0,
        3.8,
        4.1,
        "P1",
        ["渲染后 verify", "低风险 repair", "风格继承", "课程记忆"],
        RGBColor(5, 150, 105),
    )
    add_card(
        s,
        8.75,
        2.0,
        3.75,
        4.1,
        "当前边界",
        ["不夸大未完成能力", "不把一键生成当主叙事", "答辩优先展示已跑通闭环"],
        RGBColor(217, 119, 6),
    )
    add_footer(s, "Page 9", color=RGBColor(120, 113, 108))

    # 10 close
    s = prs.slides.add_slide(blank)
    add_bg(s, RGBColor(6, 44, 34))
    add_title(s, "FastPPT 的核心价值", y=1.95, size=46, color=RGBColor(220, 252, 231))
    add_subtitle(s, "把老师已经想清楚的那堂课，准确、快速、可修改、可积累地做出来", y=3.0, size=24, color=RGBColor(167, 243, 208))
    add_banner(s, 1.35, 4.35, 10.7, 0.85, "理解老师 > 利用资料 > 生成课件 > 稳定修改", RGBColor(187, 247, 208), RGBColor(6, 78, 59), 22)
    add_subtitle(s, "A04 赛题提交版", y=5.55, size=18, color=RGBColor(187, 247, 208))
    add_footer(s, "Page 10", color=RGBColor(134, 239, 172))

    prs.save("02_项目简介PPT.pptx")
    print("02_项目简介PPT.pptx")


if __name__ == "__main__":
    build()
